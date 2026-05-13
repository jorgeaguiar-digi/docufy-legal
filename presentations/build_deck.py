from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────
NAVY     = RGBColor(0x0A, 0x0E, 0x2E)   # dark navy (background)
GOLD     = RGBColor(0xC9, 0xA8, 0x5C)   # Docufy gold
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xE8, 0xEC, 0xF5)   # off-white for body
RED_ACC  = RGBColor(0xE8, 0x3A, 0x3A)   # accent for pain slides

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        text="", font_size=28, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, bg_color=None, transparency=0):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_line(slide, top=1.1):
    rect(slide, 0.55, top, 1.2, 0.07, GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — HOOK / APERTURA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.5, 12, 0.9,
    "EL LUNES POR LA MAÑANA.",
    font_size=44, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 0.55, 1.55, 11, 1.1,
    "Son las 9h. Abres el móvil.",
    font_size=32, bold=False, color=WHITE)

box(s, 0.55, 2.7, 11, 2.8,
    "11 WhatsApps de clientes preguntando cómo va su caso.\n"
    "Un expediente parado 10 días porque falta el parte amistoso.\n"
    "Tu administrativa rehaciendo una reclamación por un error en la matrícula.",
    font_size=24, color=LIGHT)

accent_line(s, top=6.8)
box(s, 0.55, 6.75, 11, 0.6,
    "¿Os suena?",
    font_size=22, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — PAIN 1  (mano alzada)
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)   # left accent bar

box(s, 0.55, 1.2, 12, 1.0,
    "¿TIENES AHORA MISMO ALGÚN EXPEDIENTE",
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
box(s, 0.55, 2.2, 12, 1.0,
    "PARADO PORQUE EL CLIENTE NO HA MANDADO UN DOCUMENTO?",
    font_size=34, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 0.55, 4.0, 11, 1.0,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — PAIN 2
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 1.2, 12, 1.0,
    "¿LE HAS PEDIDO EL DNI A UN CLIENTE",
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
box(s, 0.55, 2.2, 12, 1.0,
    "MÁS DE DOS VECES?",
    font_size=34, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 0.55, 4.0, 11, 1.0,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — PAIN 3
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 1.2, 12, 1.0,
    "¿TE HA LLAMADO HOY UN CLIENTE",
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
box(s, 0.55, 2.2, 12, 1.0,
    "PREGUNTANDO EN QUÉ ESTADO ESTÁ SU CASO?",
    font_size=34, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 0.55, 4.0, 11, 1.0,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — THE STATEMENT
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 1.0, 1.6, 11, 1.2,
    "El problema no es que no seáis buenos abogados.",
    font_size=30, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 3.1, 6.3, 0.07, GOLD)

box(s, 0.8, 3.4, 11.7, 1.5,
    "Es que la mitad de vuestro tiempo\nno lo estáis ejerciendo.",
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — HOY vs EN 3 MESES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.3, 12, 0.7,
    "¿CÓMO PODRÍA SER TU SEMANA?",
    font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Left column — HOY
rect(s, 0.4, 1.15, 5.7, 5.9, RGBColor(0x14, 0x1A, 0x42))
box(s, 0.55, 1.2, 5.4, 0.65,
    "HOY", font_size=26, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
rect(s, 0.4, 1.85, 5.7, 0.05, RED_ACC)

pains = [
    "Persigues documentos por WhatsApp",
    "Copias datos a mano en Word",
    "No sabes en qué estado está cada expediente",
    "Clientes que llaman constantemente",
    "Rechazas casos por falta de capacidad",
]
for i, txt in enumerate(pains):
    box(s, 0.6, 2.05 + i*0.88, 5.3, 0.78,
        f"✗  {txt}", font_size=17, color=LIGHT)

# Right column — EN 3 MESES
rect(s, 7.2, 1.15, 5.7, 5.9, RGBColor(0x0D, 0x2A, 0x1E))
box(s, 7.3, 1.2, 5.5, 0.65,
    "EN 3 MESES", font_size=26, bold=True, color=RGBColor(0x4C, 0xD9, 0x7A), align=PP_ALIGN.CENTER)
rect(s, 7.2, 1.85, 5.7, 0.05, RGBColor(0x4C, 0xD9, 0x7A))

goods = [
    "Documentación recibida automáticamente",
    "Escritos generados en segundos",
    "Panel con estado en tiempo real",
    "El cliente recibe actualizaciones solo",
    "El doble de expedientes, mismo equipo",
]
GREEN = RGBColor(0x4C, 0xD9, 0x7A)
for i, txt in enumerate(goods):
    box(s, 7.4, 2.05 + i*0.88, 5.3, 0.78,
        f"✓  {txt}", font_size=17, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — CASO REAL (Javier)
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 0.35, 11, 0.6,
    "UN CASO REAL", font_size=20, bold=True, color=GOLD)
accent_line(s, top=1.0)

box(s, 0.55, 1.15, 11, 1.1,
    "Javier. Despacho en Madrid. 12 años en tráfico.",
    font_size=28, bold=True, color=WHITE)

box(s, 0.55, 2.35, 11, 1.5,
    '"Teníamos dos administrativas a jornada completa solo persiguiendo a clientes\npara conseguir los DNI, partes y fotos. Empezábamos el lunes\ncon la bandeja limpia y el miércoles ya estábamos otra vez colapsados."',
    font_size=20, italic=True, color=LIGHT)

# 3 stat boxes
stats = [
    ("x2", "expedientes/mes\ncon el mismo equipo"),
    ("−80%", "tiempo de gestión\nadministrativa"),
    ("+35%", "facturación media\npor caso"),
]
for i, (num, label) in enumerate(stats):
    left = 0.55 + i * 4.2
    rect(s, left, 4.3, 3.8, 2.5, RGBColor(0x12, 0x18, 0x3C))
    box(s, left + 0.15, 4.45, 3.5, 1.0,
        num, font_size=46, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    box(s, left + 0.15, 5.5, 3.5, 1.1,
        label, font_size=17, color=LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — SOLUCIÓN INTRO
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.6, 12, 0.9,
    "¿CÓMO LO HACEMOS?",
    font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
accent_line(s, top=1.6)

# 3 pillar icons / labels
pillars = [
    ("01", "Captación\ndocumental\nautomática"),
    ("02", "Organización e\ninterpretación\ncon IA"),
    ("03", "Generación de\nescritos legales\nautomatizada"),
]
for i, (num, label) in enumerate(pillars):
    left = 0.9 + i * 4.1
    rect(s, left, 2.0, 3.5, 4.8, RGBColor(0x12, 0x18, 0x3C))
    rect(s, left, 2.0, 3.5, 0.08, GOLD)
    box(s, left + 0.2, 2.15, 3.1, 0.9,
        num, font_size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    box(s, left + 0.1, 3.2, 3.3, 1.8,
        label, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — PILAR 1
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 0.3, 4.0, 0.6,
    "01 / CAPTACIÓN DOCUMENTAL", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.5,
    "Nunca más perseguiréis\nun documento.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.85, 11, 2.6,
    "El sistema contacta automáticamente con el cliente, le pide cada documento, "
    "valida que esté completo y correcto, y organiza todo en el expediente.\n\n"
    "Sin llamadas. Sin WhatsApps manuales. Sin recordatorios.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.7, 11.5, 0.08, GOLD)
box(s, 0.55, 5.85, 11, 0.8,
    "80% de clientes entregan documentación completa en menos de 48h.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 10 — PILAR 2
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 0.3, 6.0, 0.6,
    "02 / ORGANIZACIÓN E INTERPRETACIÓN CON IA", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.5,
    "El sistema lee los documentos.\nVosotros los usáis.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.85, 11, 2.6,
    "La IA extrae automáticamente los datos de cada documento: "
    "nombre, DNI, matrícula, fecha del accidente, lesiones, importes.\n\n"
    "Nada de copiar. Nada de transcribir. Nada de errores.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.7, 11.5, 0.08, GOLD)
box(s, 0.55, 5.85, 11, 0.8,
    "Eliminación del 90% de errores humanos en la documentación enviada a aseguradoras.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 11 — PILAR 3
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 0.3, 6.5, 0.6,
    "03 / GENERACIÓN DE ESCRITOS LEGALES", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.5,
    "Vuestra reclamación previa,\ngenerada en segundos.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.85, 11, 2.6,
    "El sistema genera automáticamente todos los escritos legales "
    "usando vuestras propias plantillas. "
    "Reclamaciones previas, requerimientos, minutas.\n\n"
    "Revisáis. Firmáis. Enviáis. Solo eso.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.7, 11.5, 0.08, GOLD)
box(s, 0.55, 5.85, 11, 0.8,
    "De 30 minutos por escrito a menos de 10 segundos.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 12 — RESULTADO GLOBAL
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.35, 12, 0.7,
    "EL RESULTADO",
    font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
accent_line(s, top=1.1)

metrics = [
    ("−80%",  "tiempo administrativo\npor expediente"),
    ("×2",    "capacidad operativa\nmismo equipo"),
    ("0",     "documentos perdidos\no mal cumplimentados"),
    ("+35%",  "facturación media\npor caso"),
]
for i, (num, label) in enumerate(metrics):
    col = i % 2
    row = i // 2
    left = 0.8 + col * 6.3
    top  = 1.5 + row * 2.8
    rect(s, left, top, 5.8, 2.5, RGBColor(0x12, 0x18, 0x3C))
    rect(s, left, top, 5.8, 0.08, GOLD)
    box(s, left + 0.2, top + 0.15, 5.4, 1.1,
        num, font_size=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    box(s, left + 0.2, top + 1.3, 5.4, 1.0,
        label, font_size=19, color=LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 13 — OFERTA ADEVI
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 13.33, 0.12, GOLD)   # top gold bar
rect(s, 0, 7.38, 13.33, 0.12, GOLD) # bottom gold bar

box(s, 0.55, 0.4, 12, 0.8,
    "OFERTA EXCLUSIVA PARA ASOCIADOS ADEVI",
    font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

box(s, 0.55, 1.3, 12, 1.0,
    "Solo para los presentes hoy. No está en nuestra web.",
    font_size=22, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# Offer box
rect(s, 2.0, 2.5, 9.3, 3.2, RGBColor(0x12, 0x18, 0x3C))
rect(s, 2.0, 2.5, 9.3, 0.1, GOLD)
box(s, 2.2, 2.65, 8.9, 0.9,
    "Sesión de diagnóstico gratuita (30 min)",
    font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 2.2, 3.55, 8.9, 0.7,
    "Analizamos vuestro despacho y os decimos exactamente",
    font_size=18, color=LIGHT, align=PP_ALIGN.CENTER)
box(s, 2.2, 4.15, 8.9, 0.7,
    "cuántas horas podríais recuperar al mes.",
    font_size=18, color=LIGHT, align=PP_ALIGN.CENTER)
box(s, 2.2, 4.8, 8.9, 0.65,
    "Sin compromiso. Sin presión.",
    font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

box(s, 0.55, 5.9, 12, 0.6,
    "Escaneáis el QR  →  Reserváis vuestra sesión  →  Nosotros hacemos el resto.",
    font_size=19, color=LIGHT, align=PP_ALIGN.CENTER)

# QR placeholder
rect(s, 5.67, 6.55, 2.0, 0.7, RGBColor(0x1E, 0x28, 0x55))
box(s, 5.67, 6.55, 2.0, 0.7,
    "[ QR aquí ]", font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 14 — CIERRE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.8, 1.5, 11.7, 1.3,
    "Acabáis de escuchar el futuro de la IA en el derecho.",
    font_size=28, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 3.1, 6.3, 0.08, GOLD)

box(s, 0.8, 3.3, 11.7, 1.5,
    "Esto no es el futuro.\nEs lo que ya están haciendo despachos como el vuestro.",
    font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.8, 5.2, 11.7, 0.9,
    "La pregunta es cuándo queréis empezar vosotros.",
    font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, 0, 6.7, 13.33, 0.8, RGBColor(0x12, 0x18, 0x3C))
box(s, 0.5, 6.72, 12.3, 0.65,
    "docufylegal.com  ·  info@docufylegal.com  ·  +34 681 801 356",
    font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────
out = r"c:\Users\Jorge\CLAUDE\docufy-legal\presentations\ADEVI_2026_Docufy_Legal.pptx"
prs.save(out)
print("Saved: " + out)
