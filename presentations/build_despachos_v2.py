import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────
LOGO = r"c:\Users\Jorge\CLAUDE\docufy-legal\brand\Docufy Legal - logo (3).png"

# ── Brand colours ──────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x0E, 0x2E)
GOLD    = RGBColor(0xC9, 0xA8, 0x5C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE8, 0xEC, 0xF5)
CARD    = RGBColor(0x12, 0x18, 0x3C)
CARD2   = RGBColor(0x1C, 0x24, 0x58)
RED_ACC = RGBColor(0xE8, 0x3A, 0x3A)
GREEN   = RGBColor(0x4C, 0xD9, 0x7A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        text="", font_size=28, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, bg_color=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def oval(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        9, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_line(slide, top=1.1, left=0.55, width=1.4):
    rect(slide, left, top, width, 0.07, GOLD)

def left_bar(slide):
    rect(slide, 0, 0, 0.12, 7.5, GOLD)

def gold_border(slide):
    rect(slide, 0, 0, 13.33, 0.08, GOLD)
    rect(slide, 0, 7.42, 13.33, 0.08, GOLD)

def logo_corner(slide, width=1.55):
    if os.path.exists(LOGO):
        left = 13.33 - 0.3 - width
        slide.shapes.add_picture(LOGO, Inches(left), Inches(0.1), width=Inches(width))

def logo_center(slide, top=1.2, width=3.8):
    if os.path.exists(LOGO):
        left = (13.33 - width) / 2
        slide.shapes.add_picture(LOGO, Inches(left), Inches(top), width=Inches(width))

def stat_circle(slide, left, top, diameter, number, label,
                circle_color=GOLD, num_color=NAVY, label_color=LIGHT):
    oval(slide, left, top, diameter, diameter, circle_color)
    box(slide, left, top + diameter * 0.18, diameter, diameter * 0.45,
        number, font_size=58, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    box(slide, left - 0.3, top + diameter + 0.1, diameter + 0.6, 0.75,
        label, font_size=14, color=label_color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — HOOK / APERTURA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

# Decorative right panel
rect(s, 9.2, 0, 4.13, 7.5, CARD)
rect(s, 9.2, 0, 0.06, 7.5, GOLD)

# Logo inside right panel
logo_corner(s, width=2.4)

# Left content area
rect(s, 0, 0, 0.12, 7.5, GOLD)

box(s, 0.55, 0.55, 8.4, 1.1,
    "EL LUNES POR LA MAÑANA.",
    font_size=50, bold=True, color=GOLD)

box(s, 0.55, 1.75, 8.2, 0.75,
    "Son las 9h. Abres el móvil.",
    font_size=26, color=WHITE)

rect(s, 0.55, 2.6, 7.9, 0.05, GOLD)

pains = [
    "11 WhatsApps de clientes preguntando cómo va su caso.",
    "Un expediente parado 10 días porque falta el parte amistoso.",
    "Tu administrativa rehaciendo una reclamación por un error en la matrícula.",
]
for i, txt in enumerate(pains):
    rect(s, 0.55, 2.75 + i * 1.15, 7.9, 0.95, CARD2)
    rect(s, 0.55, 2.75 + i * 1.15, 0.06, 0.95, GOLD)
    box(s, 0.75, 2.82 + i * 1.15, 7.6, 0.8,
        txt, font_size=18, color=LIGHT)

box(s, 0.55, 6.85, 8.4, 0.55,
    "¿Os suena?",
    font_size=22, bold=True, color=GOLD)

# Right panel: big decorative number
box(s, 9.3, 1.8, 3.8, 3.8,
    "?", font_size=220, bold=True, color=RGBColor(0x1A, 0x22, 0x50),
    align=PP_ALIGN.CENTER)
box(s, 9.4, 5.5, 3.6, 0.9,
    "Tus clientes también\nse hacen esa pregunta.",
    font_size=17, italic=True, color=RGBColor(0x5A, 0x6A, 0x9E),
    align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — PREGUNTA 1
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.8, 12.5, 1.1,
    "¿TIENES AHORA MISMO ALGÚN EXPEDIENTE",
    font_size=40, bold=True, color=WHITE)

box(s, 0.55, 1.95, 12.5, 1.5,
    "PARADO PORQUE EL CLIENTE\nNO HA MANDADO UN DOCUMENTO?",
    font_size=44, bold=True, color=GOLD)

rect(s, 0.55, 3.7, 12.3, 0.06, RGBColor(0x2A, 0x34, 0x60))

box(s, 0.55, 6.3, 12.5, 0.7,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — PREGUNTA 2
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.8, 12.5, 1.1,
    "¿LE HAS PEDIDO EL DNI A UN CLIENTE",
    font_size=40, bold=True, color=WHITE)

box(s, 0.55, 1.95, 12.5, 1.5,
    "MÁS DE DOS VECES?",
    font_size=56, bold=True, color=GOLD)

rect(s, 0.55, 3.7, 12.3, 0.06, RGBColor(0x2A, 0x34, 0x60))

box(s, 0.55, 6.3, 12.5, 0.7,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — PREGUNTA 3
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.8, 12.5, 1.1,
    "¿TE HA LLAMADO HOY UN CLIENTE",
    font_size=40, bold=True, color=WHITE)

box(s, 0.55, 1.95, 12.5, 1.8,
    "PREGUNTANDO EN QUÉ\nESTADO ESTÁ SU CASO?",
    font_size=46, bold=True, color=GOLD)

rect(s, 0.55, 3.95, 12.3, 0.06, RGBColor(0x2A, 0x34, 0x60))

box(s, 0.55, 6.3, 12.5, 0.7,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — REENCUADRE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

# Decorative large quote mark
box(s, 0.4, 0.1, 2.0, 2.5,
    "“", font_size=180, bold=True,
    color=RGBColor(0x14, 0x1C, 0x46), align=PP_ALIGN.LEFT)

box(s, 1.5, 1.4, 10.8, 1.1,
    "El problema no es que no seáis buenos abogados.",
    font_size=28, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.8, 2.75, 5.7, 0.08, GOLD)

box(s, 0.8, 3.0, 11.7, 2.0,
    "Es que la mitad de vuestro tiempo\nno lo estáis ejerciendo.",
    font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom decorative closing quote
box(s, 11.3, 4.5, 2.0, 2.0,
    "”", font_size=180, bold=True,
    color=RGBColor(0x14, 0x1C, 0x46), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — HOY vs EN 3 MESES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.4, 0.2, 12.5, 0.75,
    "¿CÓMO PODRÍA SER TU SEMANA?",
    font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Columna HOY
rect(s, 0.3, 1.1, 5.9, 6.1, RGBColor(0x14, 0x1A, 0x42))
rect(s, 0.3, 1.1, 5.9, 0.07, RED_ACC)
box(s, 0.4, 1.2, 5.7, 0.65,
    "HOY", font_size=28, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
rect(s, 0.3, 1.85, 5.9, 0.06, RGBColor(0x30, 0x20, 0x20))

pains = [
    "Persigues documentos por WhatsApp",
    "Copias datos a mano en Word",
    "Sin visibilidad del estado de cada expediente",
    "Clientes que llaman constantemente",
    "Rechazas casos por falta de capacidad",
]
for i, txt in enumerate(pains):
    box(s, 0.45, 2.05 + i * 0.85, 5.6, 0.76,
        f"✗  {txt}", font_size=17, color=RGBColor(0xBB, 0xBB, 0xBB))

# Divider
rect(s, 6.4, 1.3, 0.5, 5.5, NAVY)
box(s, 6.3, 3.7, 0.7, 0.8,
    "→", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Columna EN 3 MESES
rect(s, 7.1, 1.1, 5.9, 6.1, RGBColor(0x0A, 0x24, 0x18))
rect(s, 7.1, 1.1, 5.9, 0.07, GREEN)
box(s, 7.2, 1.2, 5.7, 0.65,
    "EN 3 MESES", font_size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(s, 7.1, 1.85, 5.9, 0.06, RGBColor(0x10, 0x40, 0x28))

goods = [
    "Documentación recibida automáticamente",
    "Escritos generados en segundos",
    "Panel con estado en tiempo real",
    "El cliente recibe actualizaciones solo",
    "El doble de expedientes, mismo equipo",
]
for i, txt in enumerate(goods):
    box(s, 7.25, 2.05 + i * 0.85, 5.65, 0.76,
        f"✓  {txt}", font_size=17, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — LOS 3 PILARES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
logo_corner(s)

box(s, 0.55, 0.3, 11.5, 0.8,
    "¿CÓMO LO HACEMOS?",
    font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, 4.5, 1.15, 4.3, 0.07, GOLD)

# Connecting line between pillars
rect(s, 1.55, 3.55, 10.2, 0.06, RGBColor(0x2A, 0x34, 0x60))

pillars = [
    ("01", "Captación\ndocumental\nautomática",
     "Bot de WhatsApp activo\n24/7 que gestiona\ncada expediente"),
    ("02", "Organización\ne interpretación\ncon IA",
     "Extrae datos de\ncada documento\nsin intervención manual"),
    ("03", "Generación de\nescritos legales\nautomatizada",
     "Reclamaciones previas\ngeneradas en\nmenos de 10 segundos"),
]
for i, (num, title, desc) in enumerate(pillars):
    left = 0.65 + i * 4.15
    rect(s, left, 1.4, 3.7, 5.7, CARD)
    rect(s, left, 1.4, 3.7, 0.08, GOLD)
    oval(s, left + 1.1, 1.55, 1.5, 1.5, GOLD)
    box(s, left + 1.1, 1.65, 1.5, 1.3,
        num, font_size=46, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    box(s, left + 0.1, 3.2, 3.5, 1.8,
        title, font_size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, left + 0.15, 5.1, 3.4, 1.7,
        desc, font_size=14, color=RGBColor(0x8A, 0x96, 0xC4),
        align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — PILAR 1: CAPTACIÓN DOCUMENTAL
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)
logo_corner(s)

# Texto izquierda
box(s, 0.55, 0.3, 6.5, 0.6,
    "01 / CAPTACIÓN DOCUMENTAL", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.98)

box(s, 0.55, 1.15, 8.3, 1.8,
    "Nunca más perseguiréis\nun documento.",
    font_size=44, bold=True, color=WHITE)

box(s, 0.55, 3.1, 8.1, 2.4,
    "El sistema contacta automáticamente con el cliente, "
    "le pide cada documento por WhatsApp, valida que esté "
    "completo y correcto, y lo organiza en el expediente.\n\n"
    "Sin llamadas. Sin mensajes manuales. Sin recordatorios.",
    font_size=20, color=LIGHT)

# Stat circle derecha
oval(s, 9.3, 1.5, 3.3, 3.3, GOLD)
box(s, 9.3, 2.0, 3.3, 1.4,
    "80%", font_size=68, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
box(s, 9.0, 4.9, 3.9, 0.85,
    "documentación completa\nen menos de 48h",
    font_size=15, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 0.55, 5.7, 8.2, 0.07, GOLD)
box(s, 0.55, 5.85, 8.2, 0.7,
    "Se recuperan 10-15 horas semanales de trabajo administrativo.",
    font_size=17, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — PILAR 2: ORGANIZACIÓN CON IA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)
logo_corner(s)

box(s, 0.55, 0.3, 8.0, 0.6,
    "02 / ORGANIZACIÓN E INTERPRETACIÓN CON IA",
    font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.98)

box(s, 0.55, 1.15, 8.3, 1.8,
    "El sistema lee los documentos.\nVosotros los usáis.",
    font_size=44, bold=True, color=WHITE)

box(s, 0.55, 3.1, 8.1, 2.4,
    "La IA extrae automáticamente nombre, DNI, matrícula, "
    "fecha del accidente, lesiones e importes de cada documento.\n\n"
    "Nada de copiar. Nada de transcribir. Nada de errores.",
    font_size=20, color=LIGHT)

# Stat circle derecha
oval(s, 9.3, 1.5, 3.3, 3.3, GOLD)
box(s, 9.3, 1.9, 3.3, 1.6,
    "−90%", font_size=62, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
box(s, 9.0, 4.9, 3.9, 0.85,
    "errores humanos en\ndocumentación a aseguradoras",
    font_size=15, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 0.55, 5.7, 8.2, 0.07, GOLD)
box(s, 0.55, 5.85, 8.2, 0.7,
    "De 15-30 minutos extrayendo datos a 0 segundos.",
    font_size=17, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 10 — PILAR 3: GENERACIÓN DE ESCRITOS
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)
logo_corner(s)

box(s, 0.55, 0.3, 7.5, 0.6,
    "03 / GENERACIÓN DE ESCRITOS LEGALES",
    font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.98)

box(s, 0.55, 1.15, 8.3, 1.8,
    "Vuestra reclamación previa,\ngenerada en segundos.",
    font_size=44, bold=True, color=WHITE)

box(s, 0.55, 3.1, 8.1, 2.4,
    "El sistema genera automáticamente todos los escritos "
    "legales usando vuestras propias plantillas: "
    "reclamaciones previas, requerimientos, minutas.\n\n"
    "Revisáis. Firmáis. Enviáis. Solo eso.",
    font_size=20, color=LIGHT)

# Stat circle derecha
oval(s, 9.3, 1.5, 3.3, 3.3, GOLD)
box(s, 9.3, 2.05, 3.3, 1.5,
    "<10s", font_size=58, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
box(s, 9.0, 4.9, 3.9, 0.85,
    "por escrito generado\nvs 30 min de forma manual",
    font_size=15, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 0.55, 5.7, 8.2, 0.07, GOLD)
box(s, 0.55, 5.85, 8.2, 0.7,
    "Cero errores de transcripción. Mismo estilo jurídico del despacho.",
    font_size=17, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 11 — OFERTA / CTA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
gold_border(s)
logo_corner(s)

box(s, 0.6, 0.2, 11.5, 0.85,
    "SESIÓN DE DIAGNÓSTICO GRATUITA",
    font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

box(s, 1.5, 1.1, 10.3, 0.65,
    "Solo para los presentes hoy. No está en nuestra web.",
    font_size=19, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 1.5, 1.8, 10.3, 0.06, RGBColor(0x2A, 0x34, 0x60))

items = [
    ("Diagnóstico de vuestro despacho",
     "Analizamos cuellos de botella y calculamos el coste real del problema"),
    ("Demo en directo de Docufy Legal",
     "Veis el sistema funcionando con un expediente real en tiempo real"),
    ("Plan de automatización personalizado",
     "Os lleváis un roadmap concreto para los próximos 30 días"),
]
for i, (title, desc) in enumerate(items):
    top = 2.0 + i * 1.3
    rect(s, 1.5, top, 10.3, 1.1, CARD)
    rect(s, 1.5, top, 0.08, 1.1, GOLD)
    box(s, 1.75, top + 0.08, 4.5, 0.45,
        title, font_size=18, bold=True, color=WHITE)
    box(s, 1.75, top + 0.55, 9.8, 0.5,
        desc, font_size=15, color=LIGHT)

box(s, 1.5, 6.0, 8.5, 0.6,
    "Escaneáis el QR  →  Reserváis vuestra sesión  →  Nosotros hacemos el resto.",
    font_size=17, color=LIGHT)

rect(s, 10.4, 5.7, 2.2, 1.6, CARD)
box(s, 10.4, 5.7, 2.2, 1.6,
    "[ QR aquí ]", font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 12 — CIERRE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 13.33, 0.06, GOLD)

box(s, 0.8, 0.2, 11.7, 0.75,
    "Acabáis de escuchar el futuro de la IA en el derecho.",
    font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.0, 1.1, 7.3, 0.08, GOLD)

box(s, 0.6, 1.3, 12.1, 1.6,
    "Esto no es el futuro.\nEs lo que ya están haciendo despachos como el vuestro.",
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Logo grande centrado
logo_center(s, top=3.05, width=4.0)

box(s, 0.8, 5.55, 11.7, 0.7,
    "La pregunta es cuándo queréis empezar vosotros.",
    font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, 0, 6.55, 13.33, 0.95, CARD)
box(s, 0.5, 6.6, 12.3, 0.75,
    "docufylegal.com  ·  info@docufylegal.com  ·  +34 681 801 356",
    font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# ── Guardar ────────────────────────────────────────────────────────
out = r"c:\Users\Jorge\CLAUDE\docufy-legal\presentations\Docufy_Legal_v2.pptx"
prs.save(out)
print("Guardado: " + out)
