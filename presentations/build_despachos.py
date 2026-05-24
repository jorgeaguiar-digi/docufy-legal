from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paleta oficial Docufy Legal
NAVY    = RGBColor(0x0A, 0x0E, 0x2E)
GOLD    = RGBColor(0xC9, 0xA8, 0x5C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE8, 0xEC, 0xF5)
CARD    = RGBColor(0x12, 0x18, 0x3C)
RED_ACC = RGBColor(0xE8, 0x3A, 0x3A)
GREEN   = RGBColor(0x4C, 0xD9, 0x7A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        text="", font_size=28, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, bg_color=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_line(slide, top=1.1):
    rect(slide, 0.55, top, 1.2, 0.07, GOLD)

def left_bar(slide):
    rect(slide, 0, 0, 0.12, 7.5, GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — HOOK DE APERTURA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.5, 12, 0.9,
    "EL LUNES POR LA MAÑANA.",
    font_size=44, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 0.55, 1.55, 11, 0.9,
    "Son las 9h. Abres el móvil.",
    font_size=32, color=WHITE)

box(s, 0.55, 2.65, 11, 2.8,
    "11 WhatsApps de clientes preguntando cómo va su caso.\n"
    "Un expediente parado 10 días porque falta el parte amistoso.\n"
    "Tu administrativa rehaciendo una reclamación por un error en la matrícula.",
    font_size=24, color=LIGHT)

rect(s, 0.55, 6.7, 11.5, 0.06, GOLD)
box(s, 0.55, 6.75, 11, 0.6,
    "¿Os suena?",
    font_size=22, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — PREGUNTA 1
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 1.3, 12, 1.0,
    "¿TIENES AHORA MISMO ALGÚN EXPEDIENTE",
    font_size=34, bold=True, color=WHITE)
box(s, 0.55, 2.35, 12, 1.0,
    "PARADO PORQUE EL CLIENTE NO HA MANDADO UN DOCUMENTO?",
    font_size=34, bold=True, color=GOLD)

box(s, 0.55, 4.2, 11, 0.9,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — PREGUNTA 2
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 1.3, 12, 1.0,
    "¿LE HAS PEDIDO EL DNI A UN CLIENTE",
    font_size=34, bold=True, color=WHITE)
box(s, 0.55, 2.35, 12, 1.0,
    "MÁS DE DOS VECES?",
    font_size=34, bold=True, color=GOLD)

box(s, 0.55, 4.2, 11, 0.9,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — PREGUNTA 3
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 1.3, 12, 1.0,
    "¿TE HA LLAMADO HOY UN CLIENTE",
    font_size=34, bold=True, color=WHITE)
box(s, 0.55, 2.35, 12, 1.0,
    "PREGUNTANDO EN QUÉ ESTADO ESTÁ SU CASO?",
    font_size=34, bold=True, color=GOLD)

box(s, 0.55, 4.2, 11, 0.9,
    "Levantad la mano.",
    font_size=26, italic=True, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — REENCUADRE DEL PROBLEMA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 1.0, 1.6, 11.3, 1.1,
    "El problema no es que no seáis buenos abogados.",
    font_size=30, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 3.0, 6.3, 0.07, GOLD)

box(s, 0.8, 3.25, 11.7, 1.8,
    "Es que la mitad de vuestro tiempo\nno lo estáis ejerciendo.",
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — HOY vs EN 3 MESES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.3, 12, 0.7,
    "¿CÓMO PODRÍA SER TU SEMANA?",
    font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Columna HOY
rect(s, 0.4, 1.15, 5.7, 5.9, RGBColor(0x14, 0x1A, 0x42))
box(s, 0.55, 1.2, 5.4, 0.65,
    "HOY", font_size=26, bold=True, color=RED_ACC, align=PP_ALIGN.CENTER)
rect(s, 0.4, 1.85, 5.7, 0.05, RED_ACC)

pains = [
    "Persigues documentos por WhatsApp",
    "Copias datos a mano en Word",
    "Sin visibilidad del estado de cada expediente",
    "Clientes que llaman constantemente",
    "Rechazas casos por falta de capacidad",
]
for i, txt in enumerate(pains):
    box(s, 0.6, 2.05 + i * 0.88, 5.3, 0.78,
        f"✗  {txt}", font_size=17, color=LIGHT)

# Columna EN 3 MESES
rect(s, 7.2, 1.15, 5.7, 5.9, RGBColor(0x0D, 0x2A, 0x1E))
box(s, 7.3, 1.2, 5.5, 0.65,
    "EN 3 MESES", font_size=26, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(s, 7.2, 1.85, 5.7, 0.05, GREEN)

goods = [
    "Documentación recibida automáticamente",
    "Escritos generados en segundos",
    "Panel con estado en tiempo real",
    "El cliente recibe actualizaciones solo",
    "El doble de expedientes, mismo equipo",
]
for i, txt in enumerate(goods):
    box(s, 7.4, 2.05 + i * 0.88, 5.3, 0.78,
        f"✓  {txt}", font_size=17, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — LOS 3 PILARES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.55, 0.55, 12, 0.8,
    "¿CÓMO LO HACEMOS?",
    font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
accent_line(s, top=1.45)

pillars = [
    ("01", "Captación\ndocumental\nautomática"),
    ("02", "Organización e\ninterpretación\ncon IA"),
    ("03", "Generación de\nescritos legales\nautomatizada"),
]
for i, (num, label) in enumerate(pillars):
    left = 0.9 + i * 4.1
    rect(s, left, 1.9, 3.5, 4.9, CARD)
    rect(s, left, 1.9, 3.5, 0.08, GOLD)
    box(s, left + 0.2, 2.1, 3.1, 1.0,
        num, font_size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    box(s, left + 0.1, 3.2, 3.3, 1.9,
        label, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — PILAR 1: CAPTACIÓN DOCUMENTAL
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.3, 5.0, 0.6,
    "01 / CAPTACIÓN DOCUMENTAL", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.6,
    "Nunca más perseguiréis\nun documento.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.9, 11, 2.5,
    "El sistema contacta automáticamente con el cliente, le pide cada documento "
    "por WhatsApp, valida que esté completo y correcto, y lo organiza en el expediente.\n\n"
    "Sin llamadas. Sin mensajes manuales. Sin recordatorios.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.65, 11.5, 0.08, GOLD)
box(s, 0.55, 5.8, 11, 0.8,
    "80% de clientes entregan documentación completa en menos de 48h.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — PILAR 2: ORGANIZACIÓN CON IA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.3, 7.0, 0.6,
    "02 / ORGANIZACIÓN E INTERPRETACIÓN CON IA", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.6,
    "El sistema lee los documentos.\nVosotros los usáis.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.9, 11, 2.5,
    "La IA extrae automáticamente los datos de cada documento: "
    "nombre, DNI, matrícula, fecha del accidente, lesiones, importes.\n\n"
    "Nada de copiar. Nada de transcribir. Nada de errores.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.65, 11.5, 0.08, GOLD)
box(s, 0.55, 5.8, 11, 0.8,
    "Eliminación del 90% de errores humanos en documentación enviada a aseguradoras.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 10 — PILAR 3: GENERACIÓN DE ESCRITOS
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
left_bar(s)

box(s, 0.55, 0.3, 6.5, 0.6,
    "03 / GENERACIÓN DE ESCRITOS LEGALES", font_size=18, bold=True, color=GOLD)
accent_line(s, top=0.95)

box(s, 0.55, 1.1, 11.5, 1.6,
    "Vuestra reclamación previa,\ngenerada en segundos.",
    font_size=42, bold=True, color=WHITE)

box(s, 0.55, 2.9, 11, 2.5,
    "El sistema genera automáticamente todos los escritos legales "
    "usando vuestras propias plantillas: reclamaciones previas, requerimientos, minutas.\n\n"
    "Revisáis. Firmáis. Enviáis. Solo eso.",
    font_size=22, color=LIGHT)

rect(s, 0.55, 5.65, 11.5, 0.08, GOLD)
box(s, 0.55, 5.8, 11, 0.8,
    "De 30 minutos por escrito a menos de 10 segundos.",
    font_size=20, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 11 — OFERTA / CTA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 13.33, 0.1, GOLD)
rect(s, 0, 7.4, 13.33, 0.1, GOLD)

box(s, 0.55, 0.25, 12, 0.85,
    "SESIÓN DE DIAGNÓSTICO GRATUITA",
    font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

box(s, 0.55, 1.2, 12, 0.65,
    "Solo para los presentes hoy. Sin compromiso. Sin presión.",
    font_size=20, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

items = [
    "Analizamos vuestro despacho y detectamos cuellos de botella",
    "Os mostramos cómo funciona Docufy Legal en directo",
    "Recibís un plan de automatización personalizado",
]
for i, txt in enumerate(items):
    rect(s, 1.5, 2.1 + i * 1.05, 10.3, 0.82, CARD)
    box(s, 1.8, 2.18 + i * 1.05, 9.8, 0.68,
        f"✓  {txt}", font_size=19, color=LIGHT)

box(s, 0.55, 5.4, 12, 0.65,
    "Escaneáis el QR  →  Reserváis vuestra sesión  →  Nosotros hacemos el resto.",
    font_size=18, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 5.67, 6.1, 2.0, 1.1, CARD)
box(s, 5.67, 6.1, 2.0, 1.1,
    "[ QR aquí ]", font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 12 — CIERRE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

box(s, 0.8, 0.7, 11.7, 0.9,
    "Acabáis de escuchar el futuro de la IA en el derecho.",
    font_size=26, color=LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 1.8, 6.3, 0.08, GOLD)

box(s, 0.8, 2.0, 11.7, 1.5,
    "Esto no es el futuro.\nEs lo que ya están haciendo despachos como el vuestro.",
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.8, 3.7, 11.7, 1.6, CARD)
box(s, 1.0, 3.8, 11.3, 1.4,
    "Los primeros despachos que automatizaron la gestión de tráfico\n"
    "no esperaron a que fuera perfecto.\n"
    "Entraron cuando todavía se estaba construyendo. Y eso les dio ventaja.",
    font_size=19, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

box(s, 0.8, 5.55, 11.7, 0.75,
    "La pregunta es cuándo queréis empezar vosotros.",
    font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, 0, 6.55, 13.33, 0.95, CARD)
box(s, 0.5, 6.6, 12.3, 0.75,
    "docufylegal.com  ·  info@docufylegal.com  ·  +34 681 801 356",
    font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# ── Guardar ────────────────────────────────────────────────────────
out = r"c:\Users\Jorge\CLAUDE\docufy-legal\presentations\Docufy_Legal_Despachos_Trafico.pptx"
prs.save(out)
print("Guardado: " + out)
